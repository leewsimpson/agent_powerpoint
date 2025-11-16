from __future__ import annotations

import base64
import hashlib
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional

from openai import AzureOpenAI, OpenAI

from .config import OpenAIConfig
from .logging_config import get_logger, log_ai_request, log_ai_response
from .prompt_store import PromptStore
from .types import ImageInput, ScoreBreakdown

logger = get_logger(__name__)


@dataclass
class ScriptGenerationResult:
    script: str
    request_id: str
    prompt_payload: str


class OpenAIClient:
    """High level abstraction over LLM powered behaviors."""

    def __init__(self, config: OpenAIConfig, prompt_store: PromptStore | None = None) -> None:
        self._config = config
        self._prompt_store = prompt_store or PromptStore()
        self._client: Optional[OpenAI | AzureOpenAI] = None
        if not config.mock_mode and config.api_key:
            if config.use_azure:
                # Initialize Azure OpenAI client
                if not config.azure_endpoint:
                    raise ValueError("azure_endpoint is required for Azure OpenAI")
                if not config.azure_api_version:
                    raise ValueError("azure_api_version is required for Azure OpenAI")
                self._client = AzureOpenAI(
                    api_key=config.api_key,
                    azure_endpoint=config.azure_endpoint,
                    api_version=config.azure_api_version,
                )
            else:
                # Initialize standard OpenAI client
                self._client = OpenAI(api_key=config.api_key)

    def generate_initial_script(
            self, 
            prompt: str,
            reference_image: Optional[Path] = None,
            image_assets: Iterable[ImageInput] = ()) -> ScriptGenerationResult:

        prompt_payload = self._render_template(
            "initial_script",
            slide_brief=prompt,
            image_assets=self._format_images(image_assets),
        )
        
        log_ai_request(logger=logger, operation="GENERATE INITIAL SCRIPT", prompt=prompt_payload, reference_image=reference_image, model=self._config.default_model)
        
        if self._config.mock_mode or not self._client:
            logger.info("Using mock mode for script generation")
            script = self._mock_render_script(prompt=prompt, reference_image=reference_image, iteration_tag="initial")
            request_id = self._mock_request_id(prompt_payload)
        else:
            script, request_id = self._call_openai_with_vision(prompt_payload=prompt_payload, reference_image=reference_image)
        
        log_ai_response(logger, "GENERATE INITIAL SCRIPT", f"Generated {len(script)} characters of script code", request_id)
        return ScriptGenerationResult(script=script, request_id=request_id, prompt_payload=prompt_payload)

    def fix_script(
        self,
        prompt: str,
        image_assets: Iterable[ImageInput],
        failing_script: str,
        errors: list[str],
    ) -> ScriptGenerationResult:
        
        image_list = list(image_assets)
        error_log = "\n".join(errors) if errors else "No error details provided"
        prompt_payload = self._render_template(
            "fix_script",
            prompt=prompt,
            image_table=self._format_images(image_list),
            failing_script=failing_script,
            error_log=error_log,
        )
        
        log_ai_request(logger=logger, operation="FIX SCRIPT", prompt=prompt_payload, model=self._config.default_model)
        
        if self._config.mock_mode or not self._client:
            logger.info("Using mock mode for script fix")
            script = self._mock_render_script(prompt, iteration_tag="fixed")
            request_id = self._mock_request_id(prompt_payload)
        else:
            script, request_id = self._call_openai_with_vision(prompt_payload=prompt_payload)
        
        log_ai_response(logger, "FIX SCRIPT", f"Generated {len(script)} characters of fixed script code", request_id)
        return ScriptGenerationResult(script=script, request_id=request_id, prompt_payload=prompt_payload)

    def improve_script(
        self,
        prompt: str,
        image_assets: Iterable[ImageInput],
        previous_script: str,
        score_feedback: Optional[ScoreBreakdown],
        iteration_index: int,
        reference_image: Optional[Path] = None,
        previous_screenshot: Optional[Path] = None,
    ) -> ScriptGenerationResult:
        
        iteration_tag = f"improved_{iteration_index}"
        prompt_payload = self._render_template(
            "improve_script",
            prompt=prompt,
            image_table=self._format_images(image_assets),
            previous_script=previous_script,
            score_feedback=self._format_score(score_feedback),
            iteration_index=iteration_index,
            previous_screenshot=str(previous_screenshot) if previous_screenshot else "None",
        )
        
        log_ai_request(logger=logger, operation=f"IMPROVE SCRIPT (iteration {iteration_index})", prompt=prompt_payload, reference_image=reference_image,  model=self._config.default_model)
        
        if self._config.mock_mode or not self._client:
            logger.info("Using mock mode for script improvement")
            script = self._mock_render_script(prompt, reference_image=reference_image, previous_screenshot=previous_screenshot, iteration_tag=iteration_tag)
            request_id = self._mock_request_id(prompt_payload)
        else:
            script, request_id = self._call_openai_with_vision(prompt_payload, reference_image=reference_image, previous_screenshot=previous_screenshot)
        
        log_ai_response(logger, f"IMPROVE SCRIPT (iteration {iteration_index})", f"Generated {len(script)} characters of improved script code", request_id)
        return ScriptGenerationResult(script=script, request_id=request_id, prompt_payload=prompt_payload)

    def score_slide(
        self,
        prompt: str,
        images: Iterable[ImageInput],
        screenshot_path: Optional[Path],
        reference_image: Optional[Path],
    ) -> ScoreBreakdown:
        """Score a slide based on prompt, assets, and optionally a reference image.
        
        Uses Vision API to analyze the generated slide screenshot against the brief.
        """

        image_list = list(images)
        prompt_payload = self._render_template(
            "score_slide",
            prompt=prompt,
            image_table=self._format_images(image_list),
            screenshot_path=str(screenshot_path) if screenshot_path else "None",
            reference_image=str(reference_image) if reference_image else "None",
        )

        log_ai_request(logger=logger, operation="SCORE SLIDE", prompt=prompt_payload, reference_image=reference_image, model=self._config.default_model)

        if self._config.mock_mode or not self._client:
            logger.info("Scoring slide (mock mode)")
            return self._mock_score_slide(prompt, prompt_payload, image_list, screenshot_path, reference_image)
        
        # Validate screenshot exists for API scoring
        if not screenshot_path:
            raise ValueError("Screenshot path is required for slide scoring")
        
        if not screenshot_path.exists():
            raise FileNotFoundError(f"Screenshot file not found: {screenshot_path}")
        
        # Call Vision API with all relevant images
        score_data = self._call_openai_for_scoring(
            prompt_payload=prompt_payload,
            screenshot_path=screenshot_path,
            reference_image=reference_image,
            asset_images=[img.path for img in image_list],
        )
        
        log_ai_response(logger, "SCORE SLIDE", f"Received scores: {score_data.to_dict()}", request_id="scoring")
        return score_data

    def _mock_score_slide(
        self,
        prompt: str,
        prompt_payload: str,
        image_list: list[ImageInput],
        screenshot_path: Optional[Path],
        reference_image: Optional[Path],
    ) -> ScoreBreakdown:
        """Mock scoring implementation for testing without API calls."""
        logger.debug("Score prompt payload: %s", prompt_payload[:200] + "..." if len(prompt_payload) > 200 else prompt_payload)

        prompt_weight = min(len(prompt) / 500.0, 1.0)
        template_weight = min(len(prompt_payload) / 2000.0, 1.0)
        image_bonus = min(len(image_list) * 0.05, 0.25)
        screenshot_bonus = 0.15 if screenshot_path and screenshot_path.exists() else 0.0
        reference_bonus = 0.05 if reference_image else 0.0

        completeness = 60.0 + 30.0 * prompt_weight + image_bonus * 100
        content_accuracy = 55.0 + 35.0 * template_weight
        layout_match = 50.0 + image_bonus * 80 + reference_bonus * 100
        visual_quality = 50.0 + screenshot_bonus * 100

        completeness = min(completeness, 95.0)
        content_accuracy = min(content_accuracy, 92.0)
        layout_match = min(layout_match, 90.0)
        visual_quality = min(visual_quality, 88.0)

        aggregate = (completeness + content_accuracy + layout_match + visual_quality) / 4

        # Generate mock issues
        mock_issues = []
        if completeness < 80:
            mock_issues.append("Consider adding more content to fully address all points in the brief")
        if layout_match < 85:
            mock_issues.append("Layout could be improved to better match the reference design")
        if visual_quality < 85:
            mock_issues.append("Visual elements could be enhanced for better presentation quality")

        breakdown = ScoreBreakdown(
            completeness=round(completeness, 2),
            content_accuracy=round(content_accuracy, 2),
            layout_match=round(layout_match, 2),
            visual_quality=round(visual_quality, 2),
            aggregate=round(aggregate, 2),
            issues=mock_issues,
        )
        
        logger.info("Score breakdown: %s", breakdown.to_dict())
        return breakdown

    def _call_openai_with_vision(
        self,
        prompt_payload: str,
        reference_image: Optional[Path] = None,
        previous_screenshot: Optional[Path] = None,
    ) -> tuple[str, str]:
        """Call OpenAI Vision API with text prompt and optional images.
        
        Supports both standard chat models (with temperature) and reasoning models
        (o1, o3) which use reasoning_effort instead.
        
        Args:
            prompt_payload: The text prompt to send
            reference_image: Optional reference image to match
            previous_screenshot: Optional screenshot from previous iteration
            
        Returns:
            Tuple of (generated_script, request_id)
            
        Raises:
            ValueError: If OpenAI client not initialized
            Exception: If API call fails
        """
        if not self._client:
            raise ValueError("OpenAI client not initialized")
        
        logger.info("Calling OpenAI Vision API with model: %s", self._config.default_model)
        
        try:
            # Build message content with text and images
            content: list[dict[str, object]] = [{"type": "text", "text": prompt_payload}]
            
            # Add reference image if provided
            if reference_image and reference_image.exists():
                logger.debug("Encoding reference image: %s (size: %d bytes)", 
                           reference_image, reference_image.stat().st_size)
                base64_image = self._encode_image(reference_image)
                mime_type = self._get_image_mime_type(reference_image)
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type};base64,{base64_image}",
                        "detail": "high"
                    }
                })
                content.append({
                    "type": "text",
                    "text": "^ This is the reference image to match."
                })
            elif reference_image:
                logger.warning("Reference image does not exist: %s", reference_image)
            
            # Add previous screenshot if provided
            if previous_screenshot and previous_screenshot.exists():
                logger.debug("Encoding previous screenshot: %s (size: %d bytes)", 
                           previous_screenshot, previous_screenshot.stat().st_size)
                base64_image = self._encode_image(previous_screenshot)
                mime_type = self._get_image_mime_type(previous_screenshot)
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type};base64,{base64_image}",
                        "detail": "high"
                    }
                })
                content.append({
                    "type": "text",
                    "text": "^ This is the previous screenshot from the last iteration."
                })
            elif previous_screenshot:
                logger.warning("Previous screenshot does not exist: %s", previous_screenshot)
            
            # Build API parameters
            # For Azure OpenAI, use deployment name if provided, otherwise use default_model
            model_name = self._config.azure_deployment if self._config.use_azure and self._config.azure_deployment else self._config.default_model
            
            api_params: dict[str, object] = {
                "model": model_name,
                "messages": [
                    {"role": "system", "content": "You are an expert Python developer. Return only executable Python code."},
                    {"role": "user", "content": content}
                ],
            }
            
            # Configure parameters based on model type
            # Check both the configured model name and the actual deployment/model being used
            is_reasoning = self._is_reasoning_model(self._config.default_model) or self._is_reasoning_model(model_name)
            
            if is_reasoning:
                # Reasoning models (o1, o3, gpt-5) don't support custom temperature
                # but do support reasoning_effort
                api_params["reasoning_effort"] = self._config.reasoning_effort
                logger.info("Using reasoning model with effort: %s", self._config.reasoning_effort)
            else:
                # Non-reasoning models support temperature
                api_params["temperature"] = 0.3
            
            response = self._client.chat.completions.create(**api_params)  # type: ignore[arg-type]
            
            # Extract response content
            if not response.choices:
                raise ValueError("No choices in API response")
            
            script = response.choices[0].message.content or ""
            request_id = response.id
            
            # Log usage statistics if available
            if hasattr(response, 'usage') and response.usage:
                logger.info(
                    "OpenAI API call successful. Request ID: %s, "
                    "Tokens: prompt=%d, completion=%d, total=%d, Response length: %d chars",
                    request_id,
                    response.usage.prompt_tokens,
                    response.usage.completion_tokens,
                    response.usage.total_tokens,
                    len(script)
                )
            else:
                logger.info("OpenAI API call successful. Request ID: %s, Response length: %d chars", 
                          request_id, len(script))
            
            # Extract code from markdown code blocks if present
            script = self._extract_code_from_markdown(script)
            
            if not script:
                logger.error("Extracted script is empty after processing")
                raise ValueError("OpenAI returned empty script")
            
            logger.info("Extracted code length: %d chars", len(script))
            return script, request_id
        except Exception as error:
            logger.error("OpenAI API call failed: %s", error, exc_info=True)
            raise

    def _call_openai_for_scoring(
        self,
        prompt_payload: str,
        screenshot_path: Path,
        reference_image: Optional[Path],
        asset_images: list[Path],
    ) -> ScoreBreakdown:
        """Call OpenAI Vision API to score a slide.
        
        Args:
            prompt_payload: The scoring prompt with criteria
            screenshot_path: Screenshot of the generated slide (required)
            reference_image: Optional reference image to compare against
            asset_images: List of user-provided image assets
            
        Returns:
            ScoreBreakdown with scores and improvement issues
            
        Raises:
            ValueError: If OpenAI client not initialized, screenshot missing, or response invalid
            FileNotFoundError: If screenshot file doesn't exist
            Exception: If API call fails
        """
        if not self._client:
            raise ValueError("OpenAI client not initialized")
        
        if not screenshot_path.exists():
            raise FileNotFoundError(f"Screenshot file not found: {screenshot_path}")
        
        logger.info("Calling OpenAI Vision API for scoring with model: %s", self._config.default_model)
        
        try:
            # Build message content with text and all relevant images
            content: list[dict[str, object]] = [{"type": "text", "text": prompt_payload}]
            
            # Add screenshot (required - the main subject to score)
            logger.debug("Encoding slide screenshot: %s", screenshot_path)
            base64_image = self._encode_image(screenshot_path)
            mime_type = self._get_image_mime_type(screenshot_path)
            content.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:{mime_type};base64,{base64_image}",
                    "detail": "high"
                }
            })
            content.append({
                "type": "text",
                "text": "^ This is the generated slide to evaluate."
            })
            
            # Add reference image if provided
            if reference_image and reference_image.exists():
                logger.debug("Encoding reference image: %s", reference_image)
                base64_image = self._encode_image(reference_image)
                mime_type = self._get_image_mime_type(reference_image)
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type};base64,{base64_image}",
                        "detail": "high"
                    }
                })
                content.append({
                    "type": "text",
                    "text": "^ This is the reference image to compare layout and style against."
                })
            
            # Build API parameters for JSON response
            model_name = self._config.azure_deployment if self._config.use_azure and self._config.azure_deployment else self._config.default_model
            
            api_params: dict[str, object] = {
                "model": model_name,
                "messages": [
                    {"role": "system", "content": "You are an expert presentation evaluator. Analyze slides objectively and return only valid JSON."},
                    {"role": "user", "content": content}
                ],
                "response_format": {"type": "json_object"},
            }
            
            # Configure parameters based on model type
            is_reasoning = self._is_reasoning_model(self._config.default_model) or self._is_reasoning_model(model_name)
            
            if is_reasoning:
                api_params["reasoning_effort"] = self._config.reasoning_effort
                logger.info("Using reasoning model for scoring with effort: %s", self._config.reasoning_effort)
            else:
                api_params["temperature"] = 0.3
            
            response = self._client.chat.completions.create(**api_params)  # type: ignore[arg-type]
            
            if not response.choices:
                raise ValueError("No choices in API response")
            
            response_text = response.choices[0].message.content or ""
            request_id = response.id
            
            logger.info("Scoring API call successful. Request ID: %s", request_id)
            
            # Parse JSON response
            import json
            score_json = json.loads(response_text)
            
            # Extract scores and issues
            completeness = float(score_json.get("completeness", 0))
            content_accuracy = float(score_json.get("content_accuracy", 0))
            layout_match = float(score_json.get("layout_match", 0))
            visual_quality = float(score_json.get("visual_quality", 0))
            issues = score_json.get("issues", [])
            
            # Ensure issues is a list of strings
            if not isinstance(issues, list):
                issues = [str(issues)]
            else:
                issues = [str(issue) for issue in issues]
            
            aggregate = (completeness + content_accuracy + layout_match + visual_quality) / 4
            
            return ScoreBreakdown(
                completeness=round(completeness, 2),
                content_accuracy=round(content_accuracy, 2),
                layout_match=round(layout_match, 2),
                visual_quality=round(visual_quality, 2),
                aggregate=round(aggregate, 2),
                issues=issues,
            )
            
        except json.JSONDecodeError as error:
            logger.error("Failed to parse scoring JSON response: %s", error)
            raise ValueError(f"Invalid JSON response from scoring API: {error}") from error
        except Exception as error:
            logger.error("Scoring API call failed: %s", error, exc_info=True)
            raise
    
    @staticmethod
    def _encode_image(image_path: Path) -> str:
        """Encode image to base64 string.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Base64 encoded string of the image
            
        Raises:
            FileNotFoundError: If image file doesn't exist
            IOError: If image cannot be read
        """
        if not image_path.exists():
            raise FileNotFoundError(f"Image file not found: {image_path}")
        
        try:
            with open(image_path, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode("utf-8")
        except Exception as error:
            raise IOError(f"Failed to encode image {image_path}: {error}") from error
    
    @staticmethod
    def _get_image_mime_type(image_path: Path) -> str:
        """Get MIME type for image based on file extension.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            MIME type string (e.g., "image/png")
        """
        ext = image_path.suffix.lower()
        mime_types = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
        }
        return mime_types.get(ext, "image/png")  # Default to PNG if unknown
    
    @staticmethod
    def _is_reasoning_model(model: str) -> bool:
        """Check if the model is a reasoning model that supports reasoning_effort.
        
        Reasoning models use a different parameter set than standard models:
        - They support reasoning_effort instead of temperature
        - They include o1, o3, gpt-5, and their variants
        
        Args:
            model: The model identifier (e.g., "gpt-4o", "o1-preview", "gpt-5")
            
        Returns:
            True if the model is a reasoning model, False otherwise
        """
        # Reasoning models include o1, o3, gpt-5, and their variants
        reasoning_model_prefixes = ("o1", "o3", "gpt-5")
        return any(model.startswith(prefix) for prefix in reasoning_model_prefixes)
    
    @staticmethod
    def _extract_code_from_markdown(text: str) -> str:
        """Extract Python code from markdown code blocks.
        
        Handles both ```python and ``` code blocks. If no code blocks are found,
        returns the original text (assuming it's already pure code).
        
        Args:
            text: Response text that may contain markdown code blocks
            
        Returns:
            Extracted Python code
        """
        lines = text.split("\n")
        in_code_block = False
        code_lines: list[str] = []
        found_code_block = False
        
        for line in lines:
            stripped = line.strip()
            
            # Check for code block markers
            if stripped.startswith("```python") or stripped.startswith("```"):
                in_code_block = not in_code_block
                if not in_code_block:
                    found_code_block = True
                continue
            
            # Add lines that are inside code blocks
            if in_code_block:
                code_lines.append(line)
        
        # If we found code blocks, return the extracted code
        if found_code_block and code_lines:
            return "\n".join(code_lines).strip()
        
        # Otherwise, return original text (it's probably already pure code)
        return text.strip()

    def _mock_render_script(
        self,
        prompt: str,
        reference_image: Optional[Path] = None,
        previous_screenshot: Optional[Path] = None,
        iteration_tag: str="mock",
    ) -> str:
        prompt_lines = [line.strip() for line in prompt.splitlines() if line.strip()]
        title = prompt_lines[0] if prompt_lines else "Auto Generated Slide"
        bullet_lines = prompt_lines[1:] if len(prompt_lines) > 1 else []

        lines: list[str] = [
            f"# Auto generated script ({iteration_tag})",
            "import argparse",
            "import json",
            "from pathlib import Path",
            "from typing import Dict, Optional",
            "",
            "from pptx import Presentation",
            "from pptx.util import Inches, Pt",
            "",
            "",
            "def build_text_frame(slide, title_text, bullet_points):",
            "    left = Inches(0.6)",
            "    top = Inches(0.6)",
            "    width = Inches(12.1)",
            "    height = Inches(3.8)",
            "    textbox = slide.shapes.add_textbox(left, top, width, height)",
            "    text_frame = textbox.text_frame",
            "    text_frame.text = title_text",
            "    text_frame.paragraphs[0].font.size = Pt(40)",
            "    text_frame.paragraphs[0].font.bold = True",
            "    for bullet in bullet_points:",
            "        paragraph = text_frame.add_paragraph()",
            "        paragraph.text = bullet",
            "        paragraph.level = 1",
            "        paragraph.font.size = Pt(20)",
            "",
            "",
            "def place_images(slide, image_map):",
            "    image_specs = []",
        ]


        lines.extend(
            [
                "    if not image_specs:",
                "        return",
                "    base_left = Inches(0.5)",
                "    base_top = Inches(4.5)",
                "    spacing = Inches(0.3)",
                "    image_width = Inches(2.8)",
                "    image_height = Inches(2.2)",
                "    max_per_row = max(1, min(3, len(image_specs)))",
                "    for index, (name, path, description) in enumerate(image_specs):",
                "        if not path:",
                "            continue",
                "        try:",
                "            column = index % max_per_row",
                "            row = index // max_per_row",
                "            left = base_left + (image_width + spacing) * column",
                "            top = base_top + (image_height + spacing) * row",
                "            slide.shapes.add_picture(path, left, top, width=image_width, height=image_height)",
                "            if description:",
                "                caption_top = min(top + image_height + Inches(0.1), Inches(7.0))",
                "                caption_box = slide.shapes.add_textbox(left, caption_top, image_width, Inches(0.6))",
                "                caption_frame = caption_box.text_frame",
                "                caption_frame.text = description",
                "                caption_frame.paragraphs[0].font.size = Pt(12)",
                "                caption_frame.paragraphs[0].font.italic = True",
                "        except Exception as error:  # pylint: disable=broad-except",
                "            print(f\"Failed to place image {name}: {error}\", flush=True)",
                "",
                "",
                "def main(output_path: Path, image_map: Optional[Dict[str, str]] = None):",
                "    image_map = image_map or {}",
                "    presentation = Presentation()",
                "    presentation.slide_width = Inches(13.33)",
                "    presentation.slide_height = Inches(7.5)",
                "    slide = presentation.slides.add_slide(presentation.slide_layouts[6])",
                "    bullet_points = []",
            ]
        )

        for bullet_line in bullet_lines or ["Generated overview based on the prompt."]:
            lines.append(f"    bullet_points.append({bullet_line!r})")

        lines.extend(
            [
                f"    build_text_frame(slide, {title!r}, bullet_points)",
                "    place_images(slide, image_map)",
                "    presentation.save(output_path)",
                "",
                "",
                "def parse_args():",
                '    parser = argparse.ArgumentParser(description="Generated slide authoring script")',
                '    parser.add_argument("--output", required=True, type=Path)',
                '    parser.add_argument("--images", required=False, type=Path)',
                "    return parser.parse_args()",
                "",
                "",
                "def load_image_map(images_path: Optional[Path]) -> Dict[str, str]:",
                "    if not images_path or not images_path.exists():",
                "        return {}",
                "    with images_path.open(\"r\", encoding=\"utf-8\") as handle:",
                "        return json.load(handle)",
                "",
                "",
                "if __name__ == \"__main__\":",
                "    args = parse_args()",
                "    image_map = load_image_map(args.images)",
                "    main(args.output, image_map)",
            ]
        )

        return "\n".join(lines)

    @staticmethod
    def _mock_request_id(seed: str) -> str:
        digest = hashlib.sha256(seed.encode("utf-8")).hexdigest()
        return f"mock-{digest[:12]}"

    def _render_template(self, name: str, **context: object) -> str:
        return self._prompt_store.render(name, **context)

    @staticmethod
    def _format_images(images: Iterable[ImageInput]) -> str:
        if not images:
            return "(no images provided)"
        lines = ["- {name}: {description} ({path})".format(name=image.name, description=image.description, path=image.path)
                 for image in images]
        return "\n".join(lines)

    @staticmethod
    def _format_score(score: Optional[ScoreBreakdown]) -> str:
        if not score:
            return "No prior score available."
        
        score_text = (
            f"Completeness={score.completeness}, "
            f"Content Accuracy={score.content_accuracy}, "
            f"Layout Match={score.layout_match}, "
            f"Visual Quality={score.visual_quality}, "
            f"Aggregate={score.aggregate}"
        )
        
        if score.issues:
            issues_text = "\n\nIssues to address:\n" + "\n".join(f"- {issue}" for issue in score.issues)
            return score_text + issues_text
        
        return score_text
