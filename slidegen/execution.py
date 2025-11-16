from __future__ import annotations

import json
import shutil
import subprocess
import sys
import time
from pathlib import Path
from typing import Dict

from pptx import Presentation

from .artifacts import ArtifactManager, RunPaths
from .config import BehaviorConfig, RuntimeConfig
from .logging_config import get_logger
from .types import ExecutionResult, ScriptStatus, ScriptVersion

logger = get_logger(__name__)


class ExecutionEngine:
    def __init__(
        self,
        artifact_manager: ArtifactManager,
        run_paths: RunPaths,
        behavior: BehaviorConfig,
        runtime: RuntimeConfig,
    ) -> None:
        self._artifact_manager = artifact_manager
        self._run_paths = run_paths
        self._behavior = behavior
        self._runtime = runtime

    def execute(self, script: ScriptVersion, image_map: Dict[str, Path]) -> ExecutionResult:
        output_path = self._run_paths.outputs_dir / f"slide_{script.version_id}.pptx"
        image_map_path = self._run_paths.input_dir / f"{script.version_id}_images.json"
        with image_map_path.open("w", encoding="utf-8") as handle:
            json.dump({name: str(path) for name, path in image_map.items()}, handle)

        command = self._build_command(script.path, output_path, image_map_path)
        
        logger.info("=" * 80)
        logger.info("SCRIPT EXECUTION: %s", script.version_id)
        logger.info("-" * 80)
        logger.info("Command: %s", ' '.join(str(c) for c in command))
        logger.info("Working directory: %s", self._run_paths.base_dir)
        logger.info("Output path: %s", output_path)
        logger.info("Image map: %s", image_map_path)
        logger.info("Timeout: %d seconds", self._behavior.execution_timeout_seconds)
        logger.info("-" * 80)

        # Build execution info header for persisted logs
        exec_info_header = f"Executing script: {script.version_id}\n"
        exec_info_header += f"Command: {' '.join(str(c) for c in command)}\n"
        exec_info_header += f"Working directory: {self._run_paths.base_dir}\n"
        exec_info_header += f"Output path: {output_path}\n"
        exec_info_header += f"Image map: {image_map_path}\n"
        exec_info_header += "-" * 60 + "\n"

        start = time.perf_counter()
        stdout = exec_info_header
        stderr = ""
        return_code = None
        
        try:
            completed = subprocess.run(
                command,
                capture_output=True,
                text=True,
                timeout=self._behavior.execution_timeout_seconds,
                cwd=self._run_paths.base_dir,
                encoding="utf-8",
                errors="replace",
            )
            duration = time.perf_counter() - start
            stdout += completed.stdout or ""
            stderr = completed.stderr or ""
            return_code = completed.returncode
            
            logger.info("Script completed in %.2fs with return code %d", duration, return_code)
            if completed.stdout:
                logger.info("STDOUT:\n%s", completed.stdout)
            if stderr:
                logger.warning("STDERR:\n%s", stderr)
                
        except subprocess.TimeoutExpired as exc:
            duration = time.perf_counter() - start
            stdout += (exc.stdout.decode("utf-8", errors="replace") if isinstance(exc.stdout, bytes) else exc.stdout) or ""
            stderr = (exc.stderr.decode("utf-8", errors="replace") if isinstance(exc.stderr, bytes) else exc.stderr) or ""
            stderr += "\nExecution timed out"
            return_code = -1
            logger.error("Script execution timed out after %.2fs", duration)
            if exc.stdout:
                logger.info("STDOUT before timeout:\n%s", exc.stdout)
            if stderr:
                logger.error("STDERR before timeout:\n%s", stderr)
                
        except Exception as exc:
            duration = time.perf_counter() - start
            stderr = f"Failed to execute script: {exc}"
            return_code = -1
            logger.error("Script execution failed: %s", exc, exc_info=True)
        
        # Execution summary
        output_exists = output_path.exists()
        output_size = output_path.stat().st_size if output_exists else 0
        
        # Add execution summary to stdout for persisted logs
        stdout += f"\n{'-' * 60}\n"
        stdout += f"Execution completed in {duration:.2f}s\n"
        stdout += f"Return code: {return_code}\n"
        stdout += f"Output file exists: {output_exists}\n"
        if output_exists:
            stdout += f"Output file size: {output_size} bytes\n"
        
        logger.info("-" * 80)
        logger.info("EXECUTION SUMMARY")
        logger.info("Duration: %.2fs", duration)
        logger.info("Return code: %d", return_code)
        logger.info("Output file exists: %s", output_exists)
        if output_exists:
            logger.info("Output file size: %d bytes", output_size)
        logger.info("=" * 80)

        self._artifact_manager.persist_execution_logs(self._run_paths, script.version_id, stdout, stderr)

        success = return_code == 0 and output_exists
        
        # Add detailed error messages
        if return_code != 0:
            error_msg = f"Script exited with code {return_code}"
            stderr += f"\n\n{error_msg}"
            logger.error(error_msg)
        if not output_exists and return_code == 0:
            error_msg = f"Script completed but did not create output file: {output_path}"
            stderr += f"\n\n{error_msg}"
            logger.error(error_msg)
        
        if success:
            try:
                self._validate_presentation(output_path)
                logger.info("Presentation validated successfully")
            except Exception as validation_error:  # pylint: disable=broad-except
                success = False
                error_msg = f"Validation error: {validation_error}"
                stderr += f"\n\n{error_msg}"
                logger.error(error_msg)

        script.status = ScriptStatus.SUCCESS if success else ScriptStatus.FAILURE
        logger.info("Script status: %s", script.status.value)

        return ExecutionResult(
            success=success,
            pptx_path=output_path if success else None,
            stdout=stdout,
            stderr=stderr,
            return_code=return_code,
            duration_seconds=duration,
        )

    @staticmethod
    def _validate_presentation(pptx_path: Path) -> None:
        logger.info("Validating presentation: %s", pptx_path)
        presentation = Presentation(pptx_path)
        if len(presentation.slides) == 0:
            raise ValueError("Presentation has no slides")
        logger.info("Presentation has %d slide(s)", len(presentation.slides))

    def _build_command(self, script_path: Path, output_path: Path, image_map_path: Path) -> list[str]:
        if self._runtime.use_uv:
            uv_path = shutil.which(self._runtime.uv_executable)
            if uv_path:
                logger.info("Using uv to execute script: %s", uv_path)
                return [
                    uv_path,
                    "run",
                    "python",
                    str(script_path),
                    "--output",
                    str(output_path),
                    "--images",
                    str(image_map_path),
                ]
            if not self._runtime.allow_python_fallback:
                error_msg = f"uv executable '{self._runtime.uv_executable}' not found and fallback disabled"
                logger.error(error_msg)
                raise RuntimeError(error_msg)
            logger.warning("uv not found, falling back to Python: %s", sys.executable)

        logger.info("Using Python interpreter: %s", sys.executable)
        return [
            sys.executable,
            str(script_path),
            "--output",
            str(output_path),
            "--images",
            str(image_map_path),
        ]
